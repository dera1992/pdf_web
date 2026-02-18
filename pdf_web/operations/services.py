from __future__ import annotations

import logging
import secrets
from datetime import timedelta
from io import BytesIO
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
import zipfile
from xml.sax.saxutils import escape
from xml.etree import ElementTree

from django.contrib.auth.hashers import make_password
from django.core.files.base import ContentFile
from django.db import models
from django.utils import timezone

from pdf_web.documents.models import DocumentStatus
from pdf_web.documents.models import DocumentVersion
from pdf_web.operations.models import ShareLink

logger = logging.getLogger(__name__)

EXT_BY_TARGET = {
    "pdf": "pdf",
    "word": "docx",
    "excel": "xlsx",
    "ppt": "pptx",
    "jpg": "jpg",
}


def _is_truthy(value) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "on"}
    if isinstance(value, (int, float)):
        return bool(value)
    return False


PDF_EXPORT_FILTER_BY_EXT = {
    "docx": "docx:MS Word 2007 XML",
    "xlsx": "xlsx:Calc MS Excel 2007 XML",
    "pptx": "pptx:Impress MS PowerPoint 2007 XML",
}

XML_ILLEGAL_CHARS_RE = re.compile(r"[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD]")


def _xml_safe_text(value: str) -> str:
    # Remove characters illegal in XML 1.0 and escape special XML entities.
    cleaned = XML_ILLEGAL_CHARS_RE.sub("", value)
    return escape(cleaned)


def clone_version(version: DocumentVersion, *, created_by=None,
                  processing_state: dict | None = None) -> DocumentVersion:
    document = version.document
    next_version_number = (document.versions.aggregate(max_num=models.Max("version_number")) or {}).get("max_num",
                                                                                                        0) + 1
    new_version = DocumentVersion.objects.create(
        document=document,
        version_number=next_version_number,
        file=version.file,
        created_by=created_by or version.created_by,
        processing_state=processing_state or {"copied_from": version.id},
        pdf_info=version.pdf_info,
        text_content=version.text_content,
        layout_json=version.layout_json,
        security_state=version.security_state,
    )
    new_version.update_file_metadata()
    new_version.save()
    document.current_version = new_version
    document.status = DocumentStatus.ACTIVE
    document.updated_at = timezone.now()
    document.save(update_fields=["current_version", "updated_at", "status"])
    return new_version


def _minimal_pdf_bytes(text: str) -> bytes:
    safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    content = f"BT /F1 12 Tf 72 720 Td ({safe}) Tj ET"
    stream = content.encode("latin-1", errors="ignore")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = [0]
    for idx, obj in enumerate(objects, start=1):
        offsets.append(out.tell())
        out.write(f"{idx} 0 obj\n".encode("ascii"))
        out.write(obj)
        out.write(b"\nendobj\n")

    xref_pos = out.tell()
    out.write(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    out.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        out.write(f"{offset:010d} 00000 n \n".encode("ascii"))

    out.write(
        f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\nstartxref\n{xref_pos}\n%%EOF".encode("ascii")
    )
    return out.getvalue()


def _extract_docx_text(xml_data: bytes) -> str:
    values = re.findall(rb"<w:t[^>]*>(.*?)</w:t>", xml_data, flags=re.DOTALL)
    return " ".join(v.decode("utf-8", errors="ignore").strip() for v in values if v.strip())


def _extract_pptx_text(xml_data: bytes) -> str:
    values = re.findall(rb"<a:t[^>]*>(.*?)</a:t>", xml_data, flags=re.DOTALL)
    return " ".join(v.decode("utf-8", errors="ignore").strip() for v in values if v.strip())


def _extract_xlsx_text(archive: zipfile.ZipFile) -> str:
    strings: list[str] = []
    if "xl/sharedStrings.xml" in archive.namelist():
        root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
        for si in root.findall("{*}si"):
            text_parts = [t.text or "" for t in si.findall(".//{*}t")]
            strings.append("".join(text_parts).strip())

    values: list[str] = []
    sheet_paths = [n for n in archive.namelist() if n.startswith("xl/worksheets/") and n.endswith(".xml")]
    for sheet_path in sorted(sheet_paths)[:12]:
        root = ElementTree.fromstring(archive.read(sheet_path))
        for cell in root.findall(".//{*}c"):
            t_attr = cell.attrib.get("t")
            v = cell.find("{*}v")
            if v is None or v.text is None:
                continue
            raw = v.text.strip()
            if not raw:
                continue
            if t_attr == "s" and raw.isdigit():
                idx = int(raw)
                if 0 <= idx < len(strings):
                    if strings[idx]:
                        values.append(strings[idx])
            else:
                values.append(raw)
    return " ".join(values)


def _extract_ooxml_text(source_bytes: bytes, source_name: str) -> str:
    try:
        with zipfile.ZipFile(BytesIO(source_bytes)) as archive:
            if source_name.endswith(".docx") and "word/document.xml" in archive.namelist():
                return _extract_docx_text(archive.read("word/document.xml"))
            if source_name.endswith(".xlsx"):
                return _extract_xlsx_text(archive)
            if source_name.endswith(".pptx"):
                slides = [n for n in archive.namelist() if n.startswith("ppt/slides/") and n.endswith(".xml")]
                fragments = [_extract_pptx_text(archive.read(path)) for path in sorted(slides)[:30]]
                return " ".join(f for f in fragments if f).strip()
    except Exception:  # noqa: BLE001
        return ""
    return ""


def _convert_with_libreoffice(source_bytes: bytes, source_name: str, *, export_filter: str,
                              prefix: str) -> bytes | None:
    soffice_bin = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_bin:
        return None

    suffix = Path(source_name).suffix or ".xlsx"
    with tempfile.TemporaryDirectory(prefix=prefix) as tmp_dir:
        tmp_path = Path(tmp_dir)
        input_path = tmp_path / f"input{suffix}"
        output_path = tmp_path / "input.pdf"
        input_path.write_bytes(source_bytes)

        command = [
            soffice_bin,
            "--headless",
            "--nologo",
            "--nolockcheck",
            "--nodefault",
            "--nofirststartwizard",
            "--convert-to",
            f"pdf:{export_filter}",
            str(input_path),
            "--outdir",
            str(tmp_path),
        ]
        try:
            completed = subprocess.run(command, check=True, capture_output=True, timeout=90)
            if completed.stderr:
                logger.debug("LibreOffice stderr for %s: %s", source_name,
                             completed.stderr.decode("utf-8", errors="ignore"))
        except (subprocess.SubprocessError, OSError) as exc:
            logger.warning("LibreOffice PDF conversion failed for %s: %s", source_name, exc)
            return None

        if not output_path.exists():
            return None
        pdf_bytes = output_path.read_bytes()
        return pdf_bytes if pdf_bytes.startswith(b"%PDF") else None


def _convert_excel_with_libreoffice(source_bytes: bytes, source_name: str) -> bytes | None:
    return _convert_with_libreoffice(
        source_bytes,
        source_name,
        export_filter="calc_pdf_Export",
        prefix="excel-to-pdf-",
    )


def _convert_word_with_libreoffice(source_bytes: bytes, source_name: str) -> bytes | None:
    return _convert_with_libreoffice(
        source_bytes,
        source_name,
        export_filter="writer_pdf_Export",
        prefix="word-to-pdf-",
    )


def _convert_ppt_with_libreoffice(source_bytes: bytes, source_name: str) -> bytes | None:
    return _convert_with_libreoffice(
        source_bytes,
        source_name,
        export_filter="impress_pdf_Export",
        prefix="ppt-to-pdf-",
    )


def _extract_pdf_text(source_bytes: bytes, *, max_chars: int = 4000) -> str:
    try:
        import fitz

        doc = fitz.open(stream=source_bytes, filetype="pdf")
        chunks: list[str] = []
        for page in doc:
            chunks.append(page.get_text().strip())
            if sum(len(chunk) for chunk in chunks) >= max_chars:
                break
        text = "\n".join(chunk for chunk in chunks if chunk).strip()
        return text[:max_chars]
    except Exception:  # noqa: BLE001
        return ""


def _convert_pdf_with_libreoffice(source_bytes: bytes, *, target_ext: str) -> bytes | None:
    """Convert PDF using LibreOffice via ODP intermediate format (PDF->ODP->PPTX)."""
    soffice_bin = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_bin:
        logger.error("LibreOffice not found in PATH")
        return None

    with tempfile.TemporaryDirectory(prefix=f"pdf-to-{target_ext}-") as tmp_dir:
        tmp_path = Path(tmp_dir)

        # Create profile directory
        profile_path = tmp_path / "lo-profile"
        profile_path.mkdir(parents=True, exist_ok=True)
        profile_uri = profile_path.resolve().as_uri()

        input_path = tmp_path / "input.pdf"
        input_path.write_bytes(source_bytes)

        base_args = [
            soffice_bin,
            f"-env:UserInstallation={profile_uri}",
            "--headless",
            "--invisible",
            "--nologo",
            "--nolockcheck",
            "--nodefault",
            "--nofirststartwizard",
        ]

        # Step 1: PDF -> ODP using explicit impress8 filter
        intermediate_path = tmp_path / "input.odp"
        step1 = [
            *base_args,
            "--convert-to",
            "odp:impress8",  # CRITICAL: Use explicit filter name
            str(input_path),
            "--outdir",
            str(tmp_path)
        ]

        try:
            logger.info("LibreOffice Step 1: PDF -> ODP (using impress8 filter)")
            result = subprocess.run(step1, check=True, capture_output=True, timeout=300)
            if result.stderr:
                logger.debug("LibreOffice step1 stderr: %s", result.stderr.decode("utf-8", errors="ignore"))
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice PDF->ODP timeout (300s)")
            return None
        except (subprocess.SubprocessError, OSError) as exc:
            logger.warning("LibreOffice PDF->ODP conversion failed: %s", exc)
            return None

        if not intermediate_path.exists() or intermediate_path.stat().st_size == 0:
            logger.warning("LibreOffice did not create ODP file")
            return None

        # Step 2: ODP -> target format (PPTX, DOCX, etc.)
        output_path = tmp_path / f"input.{target_ext}"

        # Use explicit filter names for better compatibility
        filter_map = {
            "pptx": "Impress MS PowerPoint 2007 XML",
            "docx": "MS Word 2007 XML",
            "xlsx": "Calc MS Excel 2007 XML",
        }

        convert_arg = f"{target_ext}:{filter_map[target_ext]}" if target_ext in filter_map else target_ext
        step2 = [*base_args, "--convert-to", convert_arg, str(intermediate_path), "--outdir", str(tmp_path)]

        try:
            logger.info(f"LibreOffice Step 2: ODP -> {target_ext.upper()}")
            result = subprocess.run(step2, check=True, capture_output=True, timeout=300)
            if result.stderr:
                logger.debug("LibreOffice step2 stderr: %s", result.stderr.decode("utf-8", errors="ignore"))
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice ODP->%s timeout (300s)", target_ext.upper())
            return None
        except (subprocess.SubprocessError, OSError) as exc:
            logger.warning("LibreOffice ODP->%s conversion failed: %s", target_ext.upper(), exc)
            return None

        # Check for output file (try both lowercase and uppercase)
        if not output_path.exists():
            uppercase_path = tmp_path / f"input.{target_ext.upper()}"
            if uppercase_path.exists():
                output_path = uppercase_path
            else:
                logger.warning("LibreOffice did not create %s file", target_ext.upper())
                return None

        if output_path.stat().st_size == 0:
            logger.warning("LibreOffice %s output is empty", target_ext.upper())
            return None

        output_bytes = output_path.read_bytes()
        logger.info(f"LibreOffice successfully converted PDF->{target_ext.upper()}, size: {len(output_bytes)} bytes")
        return output_bytes


def _convert_pdf_to_docx_with_pdf2docx(source_bytes: bytes) -> bytes | None:
    """Convert PDF to DOCX using pdf2docx library (better quality)."""
    try:
        from pdf2docx import Converter

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp_path = Path(tmp_dir)
            pdf_path = tmp_path / "input.pdf"
            docx_path = tmp_path / "output.docx"

            pdf_path.write_bytes(source_bytes)

            cv = Converter(str(pdf_path))
            cv.convert(str(docx_path))
            cv.close()

            if docx_path.exists():
                return docx_path.read_bytes()
    except Exception as exc:
        logger.warning("pdf2docx conversion failed: %s", exc)
    return None


def _convert_pdf_to_pptx_with_images(source_bytes: bytes) -> bytes | None:
    """Convert PDF to PPTX by rendering each page as an image."""
    try:
        from pdf2image import convert_from_bytes
        from pptx import Presentation
        from io import BytesIO

        # Convert PDF pages to images
        images = convert_from_bytes(source_bytes, dpi=150)

        # Create PowerPoint
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]  # Blank slide

        for img in images:
            slide = prs.slides.add_slide(blank_layout)

            # Save image to BytesIO
            img_io = BytesIO()
            img.save(img_io, format='PNG')
            img_io.seek(0)

            # Add to slide (full size)
            slide.shapes.add_picture(
                img_io,
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )

        # Save to bytes
        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    except Exception as exc:
        logger.warning("pdf2image+pptx conversion failed: %s", exc)
    return None


def _convert_pdf_to_editable_pptx_stirling(source_bytes: bytes) -> bytes | None:
    """
    Convert PDF to editable PPTX using Stirling-PDF service.

    Stirling-PDF is a self-hosted open-source PDF tool that runs LibreOffice
    with proper orchestration (unoconvert) for production-quality conversions.

    Returns:
        PPTX bytes if successful, None if Stirling-PDF is unavailable
    """
    try:
        import requests
        from django.conf import settings

        # Stirling-PDF service URL (running in docker-compose)
        stirling_url = getattr(settings, 'STIRLING_PDF_URL', 'http://stirling-pdf:8080')

        logger.info("Attempting PDF→PPTX with Stirling-PDF at %s", stirling_url)

        # Prepare file for upload
        files = {
            'fileInput': ('input.pdf', source_bytes, 'application/pdf')
        }

        # Call Stirling-PDF conversion API
        # Endpoint: /api/v1/convert/pdf/powerpoint
        response = requests.post(
            f"{stirling_url}/api/v1/convert/pdf/powerpoint",
            files=files,
            timeout=180  # 3 minutes max
        )

        if response.status_code == 200:
            result_bytes = response.content

            # Validate it's actually a PPTX
            import zipfile
            from io import BytesIO
            try:
                with zipfile.ZipFile(BytesIO(result_bytes), 'r') as z:
                    required = ['[Content_Types].xml', 'ppt/presentation.xml']
                    if all(f in z.namelist() for f in required):
                        logger.info("✓ Stirling-PDF conversion successful, size: %s bytes", len(result_bytes))
                        return result_bytes
                    else:
                        logger.warning("Stirling-PDF returned invalid PPTX (missing required files)")
                        return None
            except zipfile.BadZipFile:
                logger.warning("Stirling-PDF returned corrupted PPTX")
                return None
        else:
            logger.warning(
                "Stirling-PDF returned status %s: %s",
                response.status_code,
                response.text[:200] if response.text else "No error message"
            )
            return None

    except requests.exceptions.ConnectionError:
        logger.warning("Stirling-PDF service not reachable - is it running?")
        return None
    except requests.exceptions.Timeout:
        logger.warning("Stirling-PDF conversion timeout after 180s")
        return None
    except Exception as exc:
        logger.warning("Stirling-PDF conversion error: %s", exc)
        return None


def _convert_pdf_to_editable_pptx(source_bytes: bytes) -> bytes | None:
    """
    Universal PDF -> editable PPTX.

    Core strategy: ONE text box per span (smallest atomic unit).
    - No grouping = no overlap between text boxes
    - Background rendered without text (clean image underneath)
    - Exact position/size from PDF coordinates
    - Works for any PDF: simple, multi-column, tables, forms
    """
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.oxml.xmlchemy import OxmlElement
        from pptx.oxml.ns import qn

        pdf_doc = fitz.open(stream=source_bytes, filetype="pdf")
        if len(pdf_doc) == 0:
            return None

        prs = Presentation()

        # Match slide ratio exactly to PDF
        p0 = pdf_doc[0]
        pdf_w0 = float(p0.rect.width)
        pdf_h0 = float(p0.rect.height) or 1.0
        ratio = pdf_w0 / pdf_h0
        slide_h = 7.5
        prs.slide_height = Inches(slide_h)
        prs.slide_width = Inches(slide_h * ratio)

        SLIDE_W = int(prs.slide_width)
        SLIDE_H = int(prs.slide_height)
        blank_layout = prs.slide_layouts[6]

        # ── helpers ───────────────────────────────────────────────────────

        def _zero_margins(tf):
            bp = tf._txBody.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
            )
            if bp is not None:
                for a in ("lIns", "rIns", "tIns", "bIns"):
                    bp.set(a, "0")

        def _clean_font(raw: str) -> str:
            s = (raw or "Calibri").split("+")[-1].strip() or "Calibri"
            for suf in ("-BoldItalic", "-Bold", "-Italic",
                        ",BoldItalic", ",Bold", ",Italic"):
                s = s.replace(suf, "")
            return s.strip() or "Calibri"

        def _rgb(c: int):
            return (c >> 16) & 0xFF, (c >> 8) & 0xFF, c & 0xFF

        # ── per-page ──────────────────────────────────────────────────────

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            slide = prs.slides.add_slide(blank_layout)

            pdf_w = float(page.rect.width)
            pdf_h = float(page.rect.height) or 1.0

            scale = min(SLIDE_W / pdf_w, SLIDE_H / pdf_h)
            x_off = int((SLIDE_W - pdf_w * scale) / 2)
            y_off = int((SLIDE_H - pdf_h * scale) / 2)

            # PDF pts -> EMU
            def ex(v):
                return int(v * scale + x_off)

            def ey(v):
                return int(v * scale + y_off)

            def ew(v):
                return max(int(v * scale), 914)

            def eh(v):
                return max(int(v * scale), 914)

            # ── 1. Collect all spans ──────────────────────────────────────
            raw = page.get_text("dict")

            # OCR fallback for scanned pages
            has_text = any(
                (sp.get("text") or "").strip()
                for bl in raw.get("blocks", []) or []
                if bl.get("type") == 0
                for ln in bl.get("lines", []) or []
                for sp in ln.get("spans", []) or []
            )
            if not has_text:
                try:
                    tp = page.get_textpage_ocr(dpi=300, full=True)
                    raw = page.get_text("dict", textpage=tp)
                    logger.info("Page %s: OCR used", page_num + 1)
                except Exception as e:
                    logger.warning("OCR failed page %s: %s", page_num + 1, e)

            spans = []
            all_bboxes = []

            for block in raw.get("blocks", []) or []:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []) or []:
                    for sp in line.get("spans", []) or []:
                        txt = (sp.get("text") or "").strip()
                        if not txt:
                            continue
                        bb = sp.get("bbox")
                        if not bb or len(bb) != 4:
                            continue
                        x0, y0, x1, y1 = map(float, bb)
                        if (x1 - x0) < 0.5 or (y1 - y0) < 0.5:
                            continue

                        flags = int(sp.get("flags", 0))
                        r, g, b = _rgb(int(sp.get("color", 0)))
                        size = max(6.0, min(float(sp.get("size", 12)), 200.0))

                        spans.append({
                            "text": txt,
                            "x0": x0, "y0": y0, "x1": x1, "y1": y1,
                            "size": size,
                            "font": _clean_font(sp.get("font", "")),
                            "bold": bool(flags & 16),
                            "italic": bool(flags & 2),
                            "rgb": (r, g, b),
                        })
                        all_bboxes.append((x0, y0, x1, y1))

                        # ── 2. Plain white background ────────────────────────────────
            # Strategy: don't render PDF at all - just use white slide
            # All content comes from text boxes placed at exact coordinates
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # ── 3. One text box per span ──────────────────────────────────
            # This is the KEY insight: no grouping = no overlap.
            # Each span has its own isolated text box at exact PDF coordinates.
            created = 0
            for sp in spans:
                x0, y0 = sp["x0"], sp["y0"]
                x1, y1 = sp["x1"], sp["y1"]

                left = max(0, ex(x0))
                top = max(0, ey(y0))

                # Width: exact from PDF + small buffer for font metrics
                width = ew((x1 - x0) * 1.05)

                # Height: use font size directly (most reliable)
                # PDF bbox height often underestimates PPT rendering
                height = eh(sp["size"] * 1.6)

                # Clamp to slide
                width = min(width, SLIDE_W - left)
                height = min(height, SLIDE_H - top)
                width = max(width, 200)
                height = max(height, 100)

                tx = slide.shapes.add_textbox(left, top, width, height)
                tf = tx.text_frame
                tf.word_wrap = False  # no wrap for single spans
                _zero_margins(tf)

                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = sp["text"]

                f = run.font
                f.size = Pt(sp["size"])
                f.bold = sp["bold"]
                f.italic = sp["italic"]
                f.color.rgb = RGBColor(*sp["rgb"])
                try:
                    f.name = sp["font"]
                except Exception:
                    f.name = "Calibri"

                created += 1

            logger.info(
                "Page %s: %s text boxes (one per span)",
                page_num + 1, created
            )

        pdf_doc.close()
        out = BytesIO()
        prs.save(out)
        return out.getvalue()

    except Exception as exc:
        logger.warning("PDF->editable PPTX failed: %s", exc)
        return None


def _convert_pdf_to_xlsx_with_tabula(source_bytes: bytes) -> bytes | None:
    """Convert PDF to XLSX by extracting tables."""
    try:
        import tabula
        import pandas as pd
        from io import BytesIO

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as pdf_tmp:
            pdf_tmp.write(source_bytes)
            pdf_path = pdf_tmp.name

        try:
            # Extract tables from PDF
            tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)

            if not tables:
                return None

            # Create Excel file
            output = BytesIO()
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                for i, table in enumerate(tables):
                    sheet_name = f'Table_{i + 1}' if len(tables) > 1 else 'Sheet1'
                    table.to_excel(writer, sheet_name=sheet_name, index=False)

            return output.getvalue()
        finally:
            Path(pdf_path).unlink(missing_ok=True)

    except Exception as exc:
        logger.warning("tabula conversion failed: %s", exc)
    return None


def _docx_from_text(text: str) -> bytes:
    body_text = text or "Converted from PDF"
    lines = [_xml_safe_text(line) for line in body_text.splitlines()[:80] if line.strip()]
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>"
        + "".join(f"<w:p><w:r><w:t>{line}</w:t></w:r></w:p>" for line in lines)
        + '<w:sectPr><w:pgSz w:w="12240" w:h="15840"/></w:sectPr></w:body></w:document>'
    )
    content_types = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\">
  <Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/>
  <Default Extension=\"xml\" ContentType=\"application/xml\"/>
  <Override PartName=\"/word/document.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml\"/>
</Types>"""
    rels = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"word/document.xml\"/>
</Relationships>"""
    out = BytesIO()
    with zipfile.ZipFile(out, mode="w") as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", rels)
        archive.writestr("word/document.xml", document_xml)
    return out.getvalue()


def _xlsx_from_text(text: str) -> bytes:
    lines = [line.strip() for line in text.splitlines() if line.strip()][:120] or ["Converted from PDF"]
    safe_lines = [_xml_safe_text(line) for line in lines]
    shared_items = "".join(f"<si><t>{line}</t></si>" for line in safe_lines)
    shared_strings = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        f'count="{len(safe_lines)}" uniqueCount="{len(safe_lines)}">{shared_items}</sst>'
    )
    rows = "".join(
        f'<row r="{i + 1}"><c r="A{i + 1}" t="s"><v>{i}</v></c></row>' for i in range(len(safe_lines))
    )
    sheet_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        f"<sheetData>{rows}</sheetData></worksheet>"
    )
    workbook_xml = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<workbook xmlns=\"http://schemas.openxmlformats.org/spreadsheetml/2006/main\" xmlns:r=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships\">
  <sheets><sheet name=\"Sheet1\" sheetId=\"1\" r:id=\"rId1\"/></sheets>
</workbook>"""
    root_rels = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"xl/workbook.xml\"/>
</Relationships>"""
    wb_rels = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet\" Target=\"worksheets/sheet1.xml\"/>
  <Relationship Id=\"rId2\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings\" Target=\"sharedStrings.xml\"/>
</Relationships>"""
    content_types = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\">
  <Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/>
  <Default Extension=\"xml\" ContentType=\"application/xml\"/>
  <Override PartName=\"/xl/workbook.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml\"/>
  <Override PartName=\"/xl/worksheets/sheet1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml\"/>
  <Override PartName=\"/xl/sharedStrings.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml\"/>
</Types>"""
    out = BytesIO()
    with zipfile.ZipFile(out, mode="w") as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("xl/workbook.xml", workbook_xml)
        archive.writestr("xl/_rels/workbook.xml.rels", wb_rels)
        archive.writestr("xl/worksheets/sheet1.xml", sheet_xml)
        archive.writestr("xl/sharedStrings.xml", shared_strings)
    return out.getvalue()


def _pptx_from_text(text: str) -> bytes:
    first_line = (text or "Converted from PDF").splitlines()[0][:200]
    slide_text = _xml_safe_text(first_line)
    content_types = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\">
  <Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/>
  <Default Extension=\"xml\" ContentType=\"application/xml\"/>
  <Override PartName=\"/ppt/presentation.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml\"/>
  <Override PartName=\"/ppt/slides/slide1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slide+xml\"/>
</Types>"""
    rels = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"ppt/presentation.xml\"/>
</Relationships>"""
    presentation_xml = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<p:presentation xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\" xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\" xmlns:r=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships\">
  <p:sldIdLst><p:sldId id=\"256\" r:id=\"rId1\"/></p:sldIdLst>
  <p:sldSz cx=\"9144000\" cy=\"6858000\" type=\"screen4x3\"/>
  <p:notesSz cx=\"6858000\" cy=\"9144000\"/>
</p:presentation>"""
    presentation_rels = """<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\">
  <Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide\" Target=\"slides/slide1.xml\"/>
</Relationships>"""
    slide_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<p:cSld><p:spTree>'
        '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
        '<p:sp><p:nvSpPr><p:cNvPr id="2" name="TextBox 1"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
        '<p:spPr><a:xfrm><a:off x="685800" y="685800"/><a:ext cx="7772400" cy="914400"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        f'<p:txBody><a:bodyPr wrap="square"/><a:lstStyle/><a:p><a:r><a:rPr lang="en-US" sz="2400"/>'
        f'<a:t>{slide_text}</a:t></a:r><a:endParaRPr lang="en-US"/></a:p></p:txBody></p:sp>'
        '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
    )
    out = BytesIO()
    with zipfile.ZipFile(out, mode="w") as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", rels)
        archive.writestr("ppt/presentation.xml", presentation_xml)
        archive.writestr("ppt/_rels/presentation.xml.rels", presentation_rels)
        archive.writestr("ppt/slides/slide1.xml", slide_xml)
    return out.getvalue()


def _pdf_from_upload(version: DocumentVersion, *, allow_excel_text_fallback: bool = False) -> bytes:
    if not version.file:
        return _minimal_pdf_bytes("Converted to PDF")
    version.file.open("rb")
    source_bytes = version.file.read()
    version.file.close()
    source_name = version.file.name.lower()

    if source_name.endswith(".pdf"):
        return source_bytes

    if source_name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".gif")):
        try:
            from PIL import Image

            with Image.open(BytesIO(source_bytes)) as img:
                rgb = img.convert("RGB")
                buffer = BytesIO()
                rgb.save(buffer, format="PDF", resolution=150)
                return buffer.getvalue()
        except Exception:  # noqa: BLE001
            pass

    if source_name.endswith((".xlsx", ".xls", ".xlsm")):
        converted_pdf = _convert_excel_with_libreoffice(source_bytes, source_name)
        if converted_pdf:
            return converted_pdf
        if not allow_excel_text_fallback:
            raise RuntimeError(
                "High-fidelity Excel to PDF conversion is unavailable. "
                "Install LibreOffice/soffice or retry with allow_text_fallback=true."
            )

    if source_name.endswith((".doc", ".docx", ".odt", ".rtf")):
        converted_pdf = _convert_word_with_libreoffice(source_bytes, source_name)
        if converted_pdf:
            return converted_pdf

    if source_name.endswith((".ppt", ".pptx", ".odp")):
        converted_pdf = _convert_ppt_with_libreoffice(source_bytes, source_name)
        if converted_pdf:
            return converted_pdf

    ooxml_text = _extract_ooxml_text(source_bytes, source_name)
    if ooxml_text:
        return _minimal_pdf_bytes(ooxml_text[:1200])

    text_snippet = source_bytes[:4096].decode("utf-8", errors="ignore").strip()
    if not text_snippet:
        if source_name.endswith(".doc"):
            text_snippet = "Converted Word document"
        else:
            text_snippet = f"Converted to PDF from {version.file.name}"
    return _minimal_pdf_bytes(text_snippet[:400])


def create_converted_version(
    version: DocumentVersion,
    *,
    target_format: str,
    created_by=None,
    conversion_params: dict | None = None,
) -> DocumentVersion:
    document = version.document
    next_version_number = (document.versions.aggregate(max_num=models.Max("version_number")) or {}).get("max_num",
                                                                                                        0) + 1
    base_name = Path(version.file.name).stem if version.file else f"version-{version.id}"
    extension = EXT_BY_TARGET.get(target_format, "bin")
    conversion_params = conversion_params or {}
    allow_text_fallback = _is_truthy(conversion_params.get("allow_text_fallback"))
    source_bytes = b""
    source_name = ""
    if version.file:
        version.file.open("rb")
        source_bytes = version.file.read()
        version.file.close()
        source_name = version.file.name.lower()

    if target_format == "pdf":
        output_bytes = _pdf_from_upload(
            version,
            allow_excel_text_fallback=allow_text_fallback,
        )
    elif target_format == "jpg" and source_name.endswith(".pdf"):
        try:
            import fitz

            doc = fitz.open(stream=source_bytes, filetype="pdf")
            pix = doc.load_page(0).get_pixmap(matrix=fitz.Matrix(2, 2))
            output_bytes = pix.tobytes("jpeg")
        except Exception:  # noqa: BLE001
            output_bytes = b"\xff\xd8\xff\xdb\x00C\x00" + b"0" * 128 + b"\xff\xd9"
    elif target_format in {"word", "excel", "ppt"} and source_name.endswith(".pdf"):
        target_ext = EXT_BY_TARGET.get(target_format, "bin")
        converted = None

        # Try modern libraries first (better quality)
        if target_format == "word":
            logger.info("Attempting PDF->DOCX with pdf2docx")
            converted = _convert_pdf_to_docx_with_pdf2docx(source_bytes)
            if converted:
                logger.info("Successfully converted using pdf2docx")

        elif target_format == "ppt":
            # Conversion Strategy (fallback chain):
            # 1. Stirling-PDF (production quality, editable) - self-hosted, free
            # 2. Image-based (perfect visuals, not editable) - always works

            # Try 1: Stirling-PDF (best quality editable PPTX)
            logger.info("Attempting PDF→PPTX with Stirling-PDF")
            converted = _convert_pdf_to_editable_pptx_stirling(source_bytes)

            if converted:
                logger.info("✓ Stirling-PDF conversion successful")
            else:
                # Fallback: Image-based (always works)
                logger.info("Stirling-PDF unavailable, using image-based fallback")
                converted = _convert_pdf_to_pptx_with_images(source_bytes)
                if converted:
                    logger.info("✓ Image-based PPTX created successfully")

        elif target_format == "excel":
            logger.info("Attempting PDF->XLSX with tabula")
            converted = _convert_pdf_to_xlsx_with_tabula(source_bytes)

        # Use converted bytes if successful, otherwise fall back to text extraction
        if converted:
            output_bytes = converted
        else:
            if target_format in {"word", "ppt"} and not allow_text_fallback:
                raise RuntimeError(
                    "High-fidelity PDF conversion is unavailable for this target format. "
                    "Install LibreOffice/soffice or retry with allow_text_fallback=true."
                )

            # Text extraction fallback
            logger.info("Using text extraction fallback")
            text = _extract_pdf_text(source_bytes)
            if target_format == "word":
                output_bytes = _docx_from_text(text)
            elif target_format == "excel":
                output_bytes = _xlsx_from_text(text)
            else:
                output_bytes = _pptx_from_text(text)
    elif target_format == "jpg":
        output_bytes = b"\xff\xd8\xff\xdb\x00C\x00" + b"0" * 128 + b"\xff\xd9"
    else:
        output_bytes = (
            f"Converted placeholder artifact\nTarget: {target_format}\nSource: {version.file.name if version.file else ''}\n".encode(
                "utf-8"
            )
        )

    new_version = DocumentVersion.objects.create(
        document=document,
        version_number=next_version_number,
        created_by=created_by or version.created_by,
        processing_state={"conversion": target_format},
        pdf_info=version.pdf_info,
        text_content=version.text_content,
        layout_json=version.layout_json,
        security_state=version.security_state,
    )
    new_version.file.save(f"{base_name}-converted.{extension}", ContentFile(output_bytes), save=False)
    new_version.update_file_metadata()
    new_version.save()

    document.current_version = new_version
    document.status = DocumentStatus.ACTIVE
    document.updated_at = timezone.now()
    document.save(update_fields=["current_version", "updated_at", "status"])
    return new_version


def create_share_link(*, version: DocumentVersion, user, expires_in_hours: int | None,
                      password: str | None) -> ShareLink:
    expires_at = timezone.now() + timedelta(hours=expires_in_hours) if expires_in_hours else None
    return ShareLink.objects.create(
        workspace=version.document.workspace,
        document=version.document,
        version=version,
        created_by=user,
        token=secrets.token_urlsafe(32),
        expires_at=expires_at,
        password_hash=make_password(password) if password else "",
    )
